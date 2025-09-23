"""Create text copies for files that are hard to read in GitHub.

The script now walks through the ``year_1`` directory that sits next to this
file. Every folder inside ``year_1`` may hold its own ``resources`` directory.
Whenever such a folder exists the script looks for PDF, DOCX, PPTX, and XLSX
files that do not yet have a matching ``.txt`` export. The new text file keeps
the full name of the source and simply adds ``.txt`` to the end (for example
``slides.pptx`` becomes ``slides.pptx.txt``).

You can still point the tool to a different starting folder by giving the path
as the first command-line argument.

The optional helpers used for each format are:

* ``pypdf`` (or ``PyPDF2``) for PDF files
* ``python-docx`` for Word documents
* ``python-pptx`` for PowerPoint decks
* ``openpyxl`` for Excel workbooks

If a helper is missing the script leaves the matching files untouched and
prints a friendly hint that explains how to install the needed package.
"""

from __future__ import annotations

import csv
import sys
from dataclasses import dataclass
from pathlib import Path
from typing import Callable, Iterable, Sequence


class MissingDependencyError(RuntimeError):
    """Raised when a converter needs an optional package that is not present."""


def ensure_directory(path: Path) -> Path:
    """Resolve ``path`` and confirm that it exists."""

    resolved = path.expanduser().resolve()
    if not resolved.exists():
        print(f"Resources directory not found: {resolved}", file=sys.stderr)
        raise SystemExit(1)
    return resolved


def list_target_files(
    root: Path,
    extensions: Iterable[str],
    *,
    only_missing_exports: bool,
) -> list[Path]:
    """Return files that end with ``extensions`` and may need a text export."""

    wanted = {ext.lower() for ext in extensions}
    targets: list[Path] = []
    for path in root.rglob("*"):
        if not path.is_file() or path.suffix.lower() not in wanted:
            continue
        export_path = path.with_name(f"{path.name}.txt")
        if only_missing_exports and export_path.exists():
            continue
        targets.append(path)
    return targets


def find_resource_directories(year_one_root: Path) -> list[Path]:
    """Return every ``resources`` directory that lives under ``year_one_root``."""

    resource_folders: list[Path] = []
    for child in sorted(year_one_root.iterdir()):
        if not child.is_dir():
            continue
        resource_dir = child / "resources"
        if resource_dir.is_dir():
            resource_folders.append(resource_dir)
    return resource_folders


def load_pdf_reader() -> type:
    """Return the PDF reader class from ``pypdf`` or ``PyPDF2``."""

    try:  # pragma: no cover - optional dependency
        from pypdf import PdfReader  # type: ignore
    except ImportError:  # pragma: no cover - optional dependency
        try:
            from PyPDF2 import PdfReader  # type: ignore
        except ImportError as error:  # pragma: no cover - optional dependency
            raise MissingDependencyError(
                "Install the 'pypdf' package to read PDF files."
            ) from error
    return PdfReader


def extract_pdf_text(reader: "PdfReader") -> str:
    """Collect text from all pages of the PDF ``reader``."""

    text_parts: list[str] = []
    for index, page in enumerate(reader.pages, start=1):
        page_text = page.extract_text() or ""
        if not page_text.strip():
            print(f"  Warning: page {index} appears to be empty.")
        text_parts.append(page_text)
    return "\n".join(text_parts)


def convert_pdf(pdf_path: Path) -> Path:
    """Convert ``pdf_path`` into a text file that keeps the PDF name."""

    PdfReader = load_pdf_reader()
    print(f"Converting {pdf_path}...")
    try:
        reader = PdfReader(str(pdf_path))
    except Exception as error:  # pragma: no cover - safety net
        raise RuntimeError(f"could not read PDF ({error})") from error

    if getattr(reader, "is_encrypted", False):
        try:
            reader.decrypt("")  # type: ignore[call-arg]
        except Exception as error:  # pragma: no cover - encryption fallback
            raise RuntimeError("PDF is encrypted and cannot be processed") from error

    text_content = extract_pdf_text(reader)
    output_path = pdf_path.with_name(f"{pdf_path.name}.txt")
    output_path.write_text(text_content, encoding="utf-8")
    return output_path


def convert_docx(docx_path: Path) -> Path:
    """Turn a Word document into a plain-text copy."""

    try:  # pragma: no cover - optional dependency
        from docx import Document  # type: ignore
    except ImportError as error:  # pragma: no cover - optional dependency
        raise MissingDependencyError(
            "Install the 'python-docx' package to read DOCX files."
        ) from error

    print(f"Converting {docx_path}...")
    document = Document(str(docx_path))
    text_parts: list[str] = []

    for paragraph in document.paragraphs:
        if paragraph.text:
            text_parts.append(paragraph.text)

    for table in document.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            text_parts.append("\t".join(cells).strip())

    output_path = docx_path.with_name(f"{docx_path.name}.txt")
    output_path.write_text("\n".join(part for part in text_parts if part), encoding="utf-8")
    return output_path


def convert_pptx(pptx_path: Path) -> Path:
    """Pull the text out of every slide in a PowerPoint deck."""

    try:  # pragma: no cover - optional dependency
        from pptx import Presentation  # type: ignore
    except ImportError as error:  # pragma: no cover - optional dependency
        raise MissingDependencyError(
            "Install the 'python-pptx' package to read PPTX files."
        ) from error

    print(f"Converting {pptx_path}...")
    presentation = Presentation(str(pptx_path))
    text_parts: list[str] = []

    for slide_number, slide in enumerate(presentation.slides, start=1):
        text_parts.append(f"# Slide {slide_number}")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                text_parts.append(shape.text)
        text_parts.append("")

    output_path = pptx_path.with_name(f"{pptx_path.name}.txt")
    output_path.write_text("\n".join(text_parts).strip(), encoding="utf-8")
    return output_path


def convert_xlsx(xlsx_path: Path) -> Path:
    """Copy Excel worksheets into a text file with tab-separated cells."""

    try:  # pragma: no cover - optional dependency
        from openpyxl import load_workbook  # type: ignore
    except ImportError as error:  # pragma: no cover - optional dependency
        raise MissingDependencyError(
            "Install the 'openpyxl' package to read XLSX files."
        ) from error

    print(f"Converting {xlsx_path}...")
    workbook = load_workbook(filename=str(xlsx_path), read_only=True, data_only=True)
    output_path = xlsx_path.with_name(f"{xlsx_path.name}.txt")

    with output_path.open("w", encoding="utf-8", newline="") as handle:
        writer = csv.writer(handle, delimiter="\t")
        for sheet in workbook.worksheets:
            handle.write(f"# Sheet {sheet.title}\n")
            for row in sheet.iter_rows(values_only=True):
                writer.writerow(["" if value is None else str(value) for value in row])
            handle.write("\n")

    return output_path


@dataclass(frozen=True)
class Converter:
    extension: str
    action: Callable[[Path], Path]


CONVERTERS: tuple[Converter, ...] = (
    Converter(".pdf", convert_pdf),
    Converter(".docx", convert_docx),
    Converter(".pptx", convert_pptx),
    Converter(".xlsx", convert_xlsx),
)


def convert_all(
    year_one_root: Path | str | None = None,
    *,
    echo: bool = True,
) -> list[str]:
    """Convert files inside ``resources`` folders and return status messages."""

    messages: list[str] = []

    if year_one_root is None:
        base_dir = ensure_directory(Path(__file__).resolve().parent / "year_1")
    else:
        base_dir = ensure_directory(Path(year_one_root))

    resource_dirs = find_resource_directories(base_dir)
    if not resource_dirs:
        note = f"No 'resources' folders found inside {base_dir}."
        if echo:
            print(note)
        return [note]

    wanted_extensions: Sequence[str] = [converter.extension for converter in CONVERTERS]

    for resource_dir in resource_dirs:
        header = f"Scanning {resource_dir}"
        if echo:
            print(header)
        messages.append(header)

        files_to_process = list_target_files(
            resource_dir, wanted_extensions, only_missing_exports=True
        )

        if not files_to_process:
            note = "  Everything is up to date."
            if echo:
                print(note)
            messages.append(note)
            continue

        for file_path in files_to_process:
            converter = next(
                converter.action
                for converter in CONVERTERS
                if converter.extension == file_path.suffix.lower()
            )
            try:
                output_path = converter(file_path)
            except MissingDependencyError as error:
                message = f"  Skipped {file_path.name}: {error}"
                if echo:
                    print(message, file=sys.stderr)
            except Exception as error:  # pragma: no cover - continue on failure
                message = f"  Error while converting {file_path.name}: {error}"
                if echo:
                    print(message, file=sys.stderr)
            else:
                message = f"  Saved text to {output_path}"
                if echo:
                    print(message)
            messages.append(message)

    return messages


def main() -> None:
    """Entry point used by the command-line interface."""

    if len(sys.argv) > 1:
        convert_all(sys.argv[1])
    else:
        convert_all()


if __name__ == "__main__":
    main()
